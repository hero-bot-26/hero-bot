# -*- coding: utf-8 -*-
"""
히어로 온·오프 세일즈 주간 리뷰 덱 생성기 (월요일 회의용)

weekly_review.main() 이 만든 컨텍스트(prev/last/inflow)를 받아 PPTX 를 그리고
구글 슬라이드로 변환 업로드한다. 숫자·차트는 전부 데이터에서 나오고, 헤드라인은
사실 서술 범위에서만 자동 조립한다(해석·액션 문구는 회의에서 다듬는 전제).

    python -m soo.hero_ops.weekly_review_deck              # 시트 + 덱 + 업로드
    python -m soo.hero_ops.weekly_review_deck --no-upload  # 로컬 pptx 만
"""
from __future__ import annotations

import re

import datetime as dt
import os
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from googleapiclient.http import MediaFileUpload

from soo.auth import build_services, get_credentials
from soo.hero_ops import weekly_review as WR

FONT = '맑은 고딕'
INK = RGBColor(0x11, 0x11, 0x11)
MUTED = RGBColor(0x8A, 0x8A, 0x8A)
LINE = RGBColor(0xDD, 0xDD, 0xDD)
BASE = RGBColor(0xCF, 0xCF, 0xCF)
SOFT = RGBColor(0xF4, 0xF4, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
UP = RGBColor(0x0F, 0x7B, 0x4F)
DOWN = RGBColor(0xC0, 0x39, 0x2B)
ACC = RGBColor(0xE4, 0x57, 0x2E)
W, H = 13.333, 7.5


def eok(v):
    if v is None:
        return '-'
    return '%.2f억' % (v / 1e8) if abs(v) >= 1e8 else format(round(v / 1e4), ',') + '만'


def man(v):
    return format(round((v or 0) / 1e4), ',') + '만'


def pct(v, sign=True):
    if v is None:
        return '-'
    return ('%+.1f%%' if sign else '%.1f%%') % (v * 100)


def wow(a, b):
    return WR.wow(a, b)


def g(d, k):
    return (d or {}).get(k)


def short(name, n=14):
    """'[시티 레저] 시어 후디드…' 처럼 앞의 대괄호 라인명은 라벨에서 뗀다."""
    s = re.sub(r'^\s*\[[^\]]*\]\s*', '', str(name or ''))
    return s[:n]


def ic(v):
    return format(int(v or 0), ',')


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = Inches(W), Inches(H)
        self.blank = self.prs.slide_layouts[6]

    def slide(self):
        return self.prs.slides.add_slide(self.blank)

    def save(self, path):
        self.prs.save(path)
        return path


def txt(sl, x, y, w, h, s, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, line=None):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, ln in enumerate(str(s).split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        r = p.add_run()
        r.text = ln
        f = r.font
        f.name, f.size, f.bold, f.color.rgb = FONT, Pt(size), bold, color
    return tb


def rect(sl, x, y, w, h, fill=None, edge=None, shape=MSO_SHAPE.RECTANGLE):
    sp = sl.shapes.add_shape(shape, Inches(x), Inches(y), Inches(max(w, 0.008)),
                             Inches(max(h, 0.008)))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if edge is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = edge
        sp.line.width = Pt(0.75)
    sp.text_frame.text = ''
    return sp


def hline(sl, x, y, w, color=LINE):
    ln = sl.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(0.75)
    return ln


def header(sl, kicker, title, note=None):
    txt(sl, 0.75, 0.48, 11.8, 0.22, kicker, 10.5, True, MUTED)
    txt(sl, 0.75, 0.75, 11.8, 0.5, title, 24, True, INK)
    hline(sl, 0.75, 1.42, 11.83)
    if note:
        txt(sl, 0.75, 1.55, 11.83, 0.24, note, 10.5, False, MUTED)


def foot(sl, s):
    hline(sl, 0.75, 6.92, 11.83)
    txt(sl, 0.75, 7.03, 11.83, 0.22, s, 8.5, False, MUTED)


# ------------------------------------------------------------------ 슬라이드
def s_cover(d, ctx, cards):
    sl = d.slide()
    rect(sl, 0, 0, W, H, INK)
    txt(sl, 0.95, 2.15, 11, 0.3, 'MUSINSA STANDARD  ·  HERO', 12, True, RGBColor(0x9A, 0x9A, 0x9A))
    txt(sl, 0.95, 2.62, 11, 0.85, '히어로 온·오프 세일즈 주간 리뷰', 38, True, WHITE)
    txt(sl, 0.95, 3.62, 11, 0.34, '%s   |   %s' % (' · '.join(ctx['deck_items']), ctx['period']),
        14, False, RGBColor(0xC8, 0xC8, 0xC8))
    hline(sl, 0.95, 4.35, 5.2, RGBColor(0x55, 0x55, 0x55))
    tiles = [(c['name'], eok(g(c['t'], 't_gmv')), '전주비 ' + pct(c['wow'])) for c in cards]
    inf = ctx['inflow'].get(ctx['deck_items'][0], {})
    m = inf.get('matched') or []
    if len(m) >= 2:
        tiles.append(('%s 온라인 유입' % ctx['deck_items'][0],
                      pct(wow(m[-2][2], m[-1][2])), '요일 맞춘 동일조건'))
    for i, (a, b, c) in enumerate(tiles[:3]):
        x = 0.95 + i * 3.72
        txt(sl, x, 4.72, 3.5, 0.22, a, 11, True, RGBColor(0x9A, 0x9A, 0x9A))
        txt(sl, x, 5.0, 3.5, 0.52, b, 29, True, WHITE)
        txt(sl, x, 5.62, 3.5, 0.24, c, 12, True, RGBColor(0x6E, 0xD0, 0xA6))
    txt(sl, 0.95, 6.72, 11, 0.22, dt.date.today().strftime('%Y.%m.%d') + ' 전략팀',
        10, False, RGBColor(0x77, 0x77, 0x77))


def s_summary(d, ctx, cards, src):
    sl = d.slide()
    header(sl, 'SUMMARY', '한 장 요약 — 주간 기준 · 전년비 → 전주비 → 목표비 순')
    for i, c in enumerate(cards[:2]):
        x = 0.75 + i * 6.2
        rect(sl, x, 1.78, 5.6, 4.95, SOFT)
        rect(sl, x, 1.78, 5.6, 0.062, INK)
        txt(sl, x + 0.42, 2.06, 4.8, 0.3, c['name'], 16, True, INK)
        txt(sl, x + 0.42, 2.44, 4.8, 0.62, eok(g(c['t'], 't_gmv')), 34, True, INK)
        for j, (k, v) in enumerate([('전년비', pct(c['yoy'])), ('전주비', pct(c['wow'])),
                                    ('목표비', pct(c['ach'], False) if c['ach'] else '미세팅')]):
            txt(sl, x + 0.42 + j * 1.6, 3.2, 1.55, 0.2, k, 9.5, True, MUTED)
            txt(sl, x + 0.42 + j * 1.6, 3.42, 1.55, 0.26, v, 13.5, True,
                INK if k == '목표비' else (UP if v.startswith('+') else
                                        (DOWN if v.startswith('-') else MUTED)))
        hline(sl, x + 0.42, 3.82, 4.76)
        for j, (k, v) in enumerate([
                ('판매수량', '%s → %s' % (ic(g(c['p'], 't_qty')), ic(g(c['t'], 't_qty')))),
                ('매총율', '%s → %s' % (pct(g(c['p'], 't_margin'), False),
                                    pct(g(c['t'], 't_margin'), False)))]):
            txt(sl, x + 0.42 + j * 2.4, 3.96, 2.3, 0.2, k, 9.5, True, MUTED)
            txt(sl, x + 0.42 + j * 2.4, 4.18, 2.3, 0.26, v, 12.5, True, INK)
        rect(sl, x + 0.42, 4.64, 4.76, 0.42, INK)
        txt(sl, x + 0.62, 4.72, 4.4, 0.26, c['msg'], 12.5, True, WHITE)
        for j, b in enumerate(c['bullets'][:3]):
            rect(sl, x + 0.45, 5.33 + j * 0.46, 0.055, 0.055, INK, shape=MSO_SHAPE.OVAL)
            txt(sl, x + 0.66, 5.24 + j * 0.46, 4.55, 0.42, b, 10.5, False, INK, line=1.25)
    foot(sl, src)


def s_split(d, ctx, c, src):
    sl = d.slide()
    inc = (g(c['t'], 't_gmv') or 0) - (g(c['p'], 't_gmv') or 0)
    on_i = (g(c['t'], 'on_gmv') or 0) - (g(c['p'], 'on_gmv') or 0)
    off_i = inc - on_i
    share = (on_i / inc) if inc else 0
    lead = '온라인' if share >= .5 else '오프라인'
    ttl = ('성장의 %.0f%%는 %s에서 나왔다' % (max(share, 1 - share) * 100, lead)) if inc > 0 \
        else '주간 증감 분해'
    header(sl, '%s · 01' % c['name'], ttl,
           '주간 GMV %s → %s (%s). 증감액 %s원의 채널별 분해'
           % (eok(g(c['p'], 't_gmv')), eok(g(c['t'], 't_gmv')), pct(c['wow']), ic(inc)))
    BX, BY = 1.05, 2.35
    mx = max(g(c['t'], 't_gmv') or 1, g(c['p'], 't_gmv') or 1)
    for i, (lb, v) in enumerate([('전주', g(c['p'], 't_gmv')), ('지난주', g(c['t'], 't_gmv'))]):
        x = BX + i * 2.05
        h = 3.15 * (v or 0) / mx
        rect(sl, x, BY + 3.15 - h, 1.35, h, INK if i else BASE)
        txt(sl, x, BY + 3.15 - h - 0.34, 1.35, 0.28, eok(v), 14, True, INK, PP_ALIGN.CENTER)
        txt(sl, x, BY + 3.26, 1.35, 0.5, lb, 10.5, False, MUTED, PP_ALIGN.CENTER)
    txt(sl, BX + 1.42, BY + 0.5, 0.6, 0.3, '→', 20, True, MUTED, PP_ALIGN.CENTER)

    CX = 5.55
    txt(sl, CX, 2.2, 6.9, 0.26, '증감액 %s원의 출처' % ic(inc), 12.5, True, INK)
    parts = [('온라인', on_i, INK, wow(g(c['p'], 'on_gmv'), g(c['t'], 'on_gmv')),
              '%s → %s개' % (ic(g(c['p'], 'on_qty')), ic(g(c['t'], 'on_qty')))),
             ('오프라인', off_i, BASE, wow(g(c['p'], 'off_gmv'), g(c['t'], 'off_gmv')),
              '%s → %s개' % (ic(g(c['p'], 'off_qty')), ic(g(c['t'], 'off_qty'))))]
    x0 = CX
    for lb, v, col, w_, q in parts:
        if inc <= 0 or v <= 0:
            continue
        ww = 6.85 * v / inc
        rect(sl, x0, 2.6, ww, 0.72, col)
        txt(sl, x0 + 0.12, 2.74, max(ww - 0.24, .3), 0.24, lb, 12, True,
            WHITE if col == INK else INK)
        txt(sl, x0 + 0.12, 2.99, max(ww - 0.24, .3), 0.22, '%.0f%%' % (v / inc * 100), 11, True,
            RGBColor(0xBB, 0xBB, 0xBB) if col == INK else MUTED)
        x0 += ww
    y = 3.62
    for lb, v, col, w_, q in parts:
        rect(sl, CX, y, 0.1, 0.52, col)
        txt(sl, CX + 0.26, y, 2.2, 0.24, lb, 12, True, INK)
        txt(sl, CX + 0.26, y + 0.26, 2.4, 0.22, q, 10, False, MUTED)
        txt(sl, CX + 2.6, y + 0.02, 2.0, 0.3, ('+' if v >= 0 else '') + man(v) + '원', 14, True, INK)
        txt(sl, CX + 4.75, y + 0.04, 1.7, 0.28, pct(w_), 13, True,
            UP if (w_ or 0) >= 0 else DOWN)
        y += 0.72
    hline(sl, CX, 5.15, 6.85)
    ps = (g(c['p'], 'on_gmv') or 0) / (g(c['p'], 't_gmv') or 1)
    ls = (g(c['t'], 'on_gmv') or 0) / (g(c['t'], 't_gmv') or 1)
    txt(sl, CX, 5.32, 6.85, 0.24, '온라인 비중 %.1f%% → %.1f%%  (%+.1f%%p)'
        % (ps * 100, ls * 100, (ls - ps) * 100), 12, True, INK)
    txt(sl, CX, 5.62, 6.85, 0.24, '매총율 %s → %s'
        % (pct(g(c['p'], 't_margin'), False), pct(g(c['t'], 't_margin'), False)), 11, False, MUTED)
    rect(sl, CX, 6.02, 6.85, 0.58, SOFT)
    both = (g(c['t'], 'on_gmv') or 0) > 0 and (g(c['t'], 'off_gmv') or 0) > 0
    offw = wow(g(c['p'], 'off_gmv'), g(c['t'], 'off_gmv'))
    note = ('오프라인도 %s로 동반 성장 — 채널 간 잠식이 아니라 총량 확대' % pct(offw)) \
        if both and (offw or 0) > 0 else ('전량 온라인 채널' if not both else
                                          '오프라인은 %s — 채널 간 이동 여부 확인 필요' % pct(offw))
    txt(sl, CX + 0.2, 6.14, 6.5, 0.34, note, 11.5, True, INK)
    foot(sl, src)


def s_sty(d, ctx, c, src):
    sl = d.slide()
    sty = [r for r in c['sty'] if (r.get('t_gmv') or r.get('t_qty'))]
    sty.sort(key=lambda r: -(r.get('t_gmv') or 0))
    sty = sty[:6]
    big = [r for r in sty if (r.get('t_gmv') or 0) > 1e7]
    top = max(big, key=lambda r: (r['_wow'] or -9)) if big else None
    header(sl, '%s · 02' % c['name'], 'STY별 주간 실적',
           '막대 = 지난주 주간 GMV · 우측 = 전주비와 온·오프 판매 구성')
    y = 2.15
    mx = max([r.get('t_gmv') or 0 for r in sty] or [1])
    for r in sty:
        bw = 5.5 * (r.get('t_gmv') or 0) / mx
        hi = top is not None and r is top and (r['_wow'] or 0) > .3
        rect(sl, 4.15, y + 0.1, bw, 0.46, ACC if hi else INK)
        txt(sl, 0.75, y + 0.1, 3.3, 0.26, (r.get('name') or '')[:24], 12, True, INK)
        txt(sl, 0.75, y + 0.36, 3.3, 0.2, r.get('sty', ''), 8.5, False, MUTED)
        txt(sl, 4.15 + bw + 0.14, y + 0.16, 1.6, 0.28, eok(r.get('t_gmv')), 12.5, True, INK)
        txt(sl, 10.35, y + 0.13, 1.3, 0.3, pct(r['_wow']), 14, True,
            ACC if hi else (UP if (r['_wow'] or 0) >= 0 else DOWN))
        onq, offq = r.get('on_qty') or 0, r.get('off_qty') or 0
        tq = onq + offq
        sx = 11.75
        if tq:
            for q, col in ((onq, INK), (offq, BASE)):
                ww = 1.35 * q / tq
                rect(sl, sx, y + 0.2, ww, 0.26, col)
                sx += ww
            txt(sl, 11.75, y + 0.5, 1.4, 0.18,
                'ON %.0f%% / OFF %.0f%%' % (onq / tq * 100, offq / tq * 100), 7.5, False, MUTED)
        y += 0.78
    txt(sl, 11.75, 1.86, 1.4, 0.2, '판매수량 구성', 8.5, True, MUTED)
    hline(sl, 0.75, 6.12, 11.83)
    notes = []
    if top is not None:
        onq, offq = top.get('on_qty') or 0, top.get('off_qty') or 0
        notes.append(((top.get('name') or '')[:26] + '  ' + pct(top['_wow']),
                      '온라인 비중 %.0f%% · 지난주 %s'
                      % (onq / max(onq + offq, 1) * 100, eok(top.get('t_gmv')))))
    if sty:
        notes.append(((sty[0].get('name') or '')[:26] + '  ' + eok(sty[0].get('t_gmv')),
                      '단일 STY 최대 볼륨 · 전주비 %s' % pct(sty[0]['_wow'])))
    for i, (a, b) in enumerate(notes[:2]):
        rect(sl, 0.75 + i * 6.0, 6.3, 0.055, 0.4, ACC if i == 0 else INK)
        txt(sl, 0.95 + i * 6.0, 6.28, 5.5, 0.22, a, 11.5, True, INK)
        txt(sl, 0.95 + i * 6.0, 6.52, 5.5, 0.34, b, 9.5, False, MUTED, line=1.2)
    foot(sl, src)


def s_inflow(d, ctx, item, c, src):
    inf = ctx['inflow'].get(item) or {}
    paths = inf.get('paths') or []
    if not paths:
        return
    sl = d.slide()
    m = inf.get('matched') or []
    tot_wow = wow(m[-2][2], m[-1][2]) if len(m) >= 2 else None
    da, db = inf['prev_days'], inf['cur_days']
    rank = sorted(paths, key=lambda p: -(wow(p['a'][0] / da, p['b'][0] / db) or -9))
    ext = next((p for p in paths if '외부' in p['path']), None)
    head = '유입 %s' % pct(tot_wow) if tot_wow is not None else '유입 동향'
    if ext:
        ew = wow(ext['a'][0] / da, ext['b'][0] / db)
        head += (' — 외부유입이 %s로 증가율 1위' % pct(ew)) if (rank and rank[0] is ext) \
            else (' — 외부유입 %s' % pct(ew))
    header(sl, '%s · 03' % item, head,
           '원천 최종 적재 %s. 경로별은 일평균 기준(전주 %d일 vs 지난주 %d일), 총량은 요일 맞춘 동일조건'
           % (inf.get('last_load', '-'), da, db))
    txt(sl, 0.75, 2.0, 4.2, 0.24, '경로별 일평균 유입 UV (%s)' % inf.get('season', ''),
        12, True, INK)
    rect(sl, 4.98, 2.09, 0.2, 0.1, BASE)
    txt(sl, 5.24, 2.03, 0.7, 0.2, '전주', 8.5, False, MUTED)
    rect(sl, 5.86, 2.09, 0.2, 0.1, INK)
    txt(sl, 6.12, 2.03, 0.9, 0.2, '지난주', 8.5, False, MUTED)
    y = 2.4
    mx = max([p['b'][0] / db for p in paths] or [1])
    for p in paths[:8]:
        a, b = p['a'][0] / da, p['b'][0] / db
        hi = ext is not None and p is ext
        rect(sl, 2.05, y + 0.03, 3.0 * a / mx, 0.13, BASE)
        rect(sl, 2.05, y + 0.19, 3.0 * b / mx, 0.17, ACC if hi else INK)
        txt(sl, 0.75, y + 0.06, 1.25, 0.28, p['path'], 11, True, ACC if hi else INK)
        wv = wow(a, b)
        txt(sl, 5.2, y + 0.07, 1.0, 0.26, pct(wv), 12, True,
            ACC if hi else (UP if (wv or 0) >= 0 else DOWN))
        txt(sl, 6.15, y + 0.09, 1.4, 0.22, '%s → %s' % (ic(a), ic(b)), 9, False, MUTED)
        y += 0.47

    RX = 7.75
    if ext:
        ea, eb = ext['a'][0] / da, ext['b'][0] / db
        ta = sum(p['a'][0] for p in paths) or 1
        tb = sum(p['b'][0] for p in paths) or 1
        rect(sl, RX, 1.95, 4.83, 1.42, INK)
        txt(sl, RX + 0.28, 2.12, 4.3, 0.24, '외부 유입', 11.5, True, RGBColor(0xAA, 0xAA, 0xAA))
        txt(sl, RX + 0.28, 2.4, 2.4, 0.44, pct(wow(ea, eb)), 26, True, WHITE)
        txt(sl, RX + 2.75, 2.5, 1.9, 0.22, '비중 %.1f%% → %.1f%%'
            % (ext['a'][0] / ta * 100, ext['b'][0] / tb * 100), 10, True,
            RGBColor(0xDD, 0xDD, 0xDD))
        txt(sl, RX + 2.75, 2.75, 1.9, 0.22, 'GMV 일평균 %s'
            % pct(wow(ext['a'][2] / da, ext['b'][2] / db)), 10, True, RGBColor(0xDD, 0xDD, 0xDD))
        txt(sl, RX + 0.28, 2.95, 4.3, 0.22, '일평균 %s → %s UV · 전환율 %s'
            % (ic(ea), ic(eb), pct(ext['b'][1] / ext['b'][0] if ext['b'][0] else None, False)),
            10, False, RGBColor(0xBB, 0xBB, 0xBB))

    daily = (inf.get('daily') or [])[-12:]
    if daily:
        txt(sl, RX, 3.6, 3.2, 0.24, '일별 유입 UV', 11.5, True, INK)
        by, bh = 3.98, 1.22
        mxd = max(v for _, v in daily) or 1
        cut = len(daily) - (m[-1][3] if m else 5)
        step = 4.78 / max(len(daily), 1)
        for i, (dd, v) in enumerate(daily):
            hgt = bh * v / mxd
            x = RX + i * step
            rect(sl, x, by + bh - hgt, min(0.3, step - 0.09), hgt, INK if i >= cut else BASE)
            txt(sl, x - 0.06, by + bh + 0.05, 0.42, 0.18, dd[5:].replace('-', '/'), 6.5,
                False, MUTED, PP_ALIGN.CENTER)
        j = max(range(1, len(daily)), key=lambda i: daily[i][1] - daily[i - 1][1])
        txt(sl, RX, 5.42, 4.83, 0.22, '%s %s → %s %s (%s)'
            % (daily[j - 1][0][5:], ic(daily[j - 1][1]), daily[j][0][5:], ic(daily[j][1]),
               pct(wow(daily[j - 1][1], daily[j][1]))), 10, True, INK)

    rect(sl, 0.75, 6.05, 11.83, 0.75, SOFT)
    gm = wow(g(c['p'], 'on_gmv'), g(c['t'], 'on_gmv'))
    pa = sum(p['a'][1] for p in paths) / max(sum(p['a'][0] for p in paths), 1)
    pb = sum(p['b'][1] for p in paths) / max(sum(p['b'][0] for p in paths), 1)
    txt(sl, 1.0, 6.16, 10.5, 0.24, '유입 %s   ·   온라인 GMV %s   ·   전환율 %.2f%% → %.2f%%'
        % (pct(tot_wow), pct(gm), pa * 100, pb * 100), 13.5, True, INK)
    drv = '트래픽' if (tot_wow or 0) >= (gm or 0) * .8 else '전환·객단가'
    txt(sl, 1.0, 6.46, 10.8, 0.24,
        '이번 온라인 성장은 %s가 주도 — 트래픽 유지·확대와 전환율 중 어디에 손댈지 결정 필요' % drv,
        10.5, False, MUTED)
    foot(sl, src)


def s_waterfall(d, ctx, c, src):
    sl = d.slide()
    pmap = {r['sty']: r for r in c['psty']}
    seg = []
    for r in c['sty']:
        pv = pmap.get(r['sty'], {}).get('t_gmv') or 0
        lv = r.get('t_gmv') or 0
        if pv == 0 and lv == 0:
            continue
        seg.append([short(r.get('name')), lv - pv, 'new' if pv == 0 else 'old'])
    seg = [x for x in seg if abs(x[1]) > 0]
    seg.sort(key=lambda x: x[1])
    news = [x for x in seg if x[2] == 'new']
    inc = (g(c['t'], 't_gmv') or 0) - (g(c['p'], 't_gmv') or 0)
    head = '성장의 전부가 신규 투입 STY' if (news and inc > 0 and
                                    sum(x[1] for x in news) >= inc) else '주간 증감 분해'
    header(sl, c['name'], head, '주간 GMV %s → %s (%s)'
           % (eok(g(c['p'], 't_gmv')), eok(g(c['t'], 't_gmv')), pct(c['wow'])))
    items = [('전주', g(c['p'], 't_gmv') or 0, 'base')] + \
            [(a, b, 'dn' if b < 0 else 'up') for a, b, _ in seg] + \
            [('지난주', g(c['t'], 't_gmv') or 0, 'base')]
    n = len(items)
    step = 11.3 / n
    bw = max(min(1.5, step - 0.35), 0.35)
    bx, by, bh = 1.0, 2.18, 2.32
    mx = (max(g(c['p'], 't_gmv') or 0, g(c['t'], 't_gmv') or 0) * 1.05) or 1
    run = 0
    for i, (lb, v, kind) in enumerate(items):
        x = bx + i * step
        if kind == 'base':
            h = bh * v / mx
            rect(sl, x, by + bh - h, bw, h, INK)
            txt(sl, x, by + bh - h - 0.32, bw, 0.26, man(v), 12, True, INK, PP_ALIGN.CENTER)
            run = v
        else:
            h = bh * abs(v) / mx
            top = by + bh - (run / mx * bh)
            if v < 0:
                rect(sl, x, top, bw, h, DOWN)
                ty = top + h + 0.06
            else:
                rect(sl, x, top - h, bw, h, UP)
                ty = top - h - 0.3
            txt(sl, x, ty, bw, 0.26, '%+d만' % round(v / 1e4), 10.5, True,
                DOWN if v < 0 else UP, PP_ALIGN.CENTER)
            run += v
        txt(sl, x - 0.14, by + bh + 0.12, bw + 0.28, 0.44, lb, 8.5, False, MUTED,
            PP_ALIGN.CENTER, line=1.2)
    hline(sl, bx, by + bh, 11.3)
    rect(sl, 0.75, 5.32, 5.72, 1.48, SOFT)
    base_wow = wow(g(c['p'], 't_gmv'), (g(c['t'], 't_gmv') or 0) - sum(x[1] for x in news))
    txt(sl, 1.0, 5.46, 5.2, 0.24, '신규 투입 STY 제외 시 %s' % pct(base_wow), 13, True,
        DOWN if (base_wow or 0) < 0 else UP)
    txt(sl, 1.0, 5.76, 5.3, 0.9,
        '신규 투입 %s원 / 기존 STY %s원\n매총율 %s → %s'
        % (man(sum(x[1] for x in news)), man(sum(x[1] for x in seg if x[2] == 'old')),
           pct(g(c['p'], 't_margin'), False), pct(g(c['t'], 't_margin'), False)),
        10.5, False, INK, line=1.3)
    pend = [r for r in c['sty'] if r.get('pending')]
    rect(sl, 6.86, 5.32, 5.72, 1.48, INK)
    txt(sl, 7.11, 5.46, 5.2, 0.24, '미발매 대기 STY %d종' % len(pend), 13, True, WHITE)
    txt(sl, 7.11, 5.76, 5.3, 0.9,
        '\n'.join((r.get('name') or '')[:32] for r in pend[:3]) or '없음',
        10.5, False, RGBColor(0xCC, 0xCC, 0xCC), line=1.3)
    foot(sl, src)


def s_action(d, ctx, cards, alerts, src):
    sl = d.slide()
    header(sl, 'ACTION', '재고 배분과 이번 주 퀵 액션',
           '재고주수 = 지난주 잔여재고 ÷ 지난주 판매수량. 채널 간 이관 판단 기준. 액션은 초안')
    txt(sl, 0.75, 2.0, 5.6, 0.24, '채널별 재고주수', 12.5, True, INK)
    rows = []
    for c in cards:
        for ch, lbl in (('on', '온라인'), ('off', '오프라인')):
            st, q = g(c['t'], 'stock_%s' % ch), g(c['t'], '%s_qty' % ch)
            if st and q:
                rows.append(('%s %s' % (c['name'], lbl), st / q, ic(st) + '개'))
    rows = rows[:4]
    y = 2.44
    mxw = max([r[1] for r in rows] or [1])
    for nm, v, sub in rows:
        w_ = 3.1 * min(v / mxw, 1)
        rect(sl, 2.35, y, w_, 0.36, ACC if v < 6 else INK)
        txt(sl, 0.75, y + 0.03, 1.55, 0.24, nm, 10, True, INK)
        txt(sl, 2.35 + w_ + 0.12, y + 0.03, 1.0, 0.26, '%.1f주' % v, 12, True, INK)
        txt(sl, 2.35, y + 0.4, 3.9, 0.2, sub, 8.5, False, MUTED)
        y += 0.72
    AX = 6.9
    txt(sl, AX, 2.0, 5.7, 0.24, '이번 주 퀵 액션 (초안 — 회의에서 확정)', 12.5, True, INK)
    yy = 2.42
    for i, (a, b) in enumerate((ctx.get('actions') or [])[:5]):
        rect(sl, AX, yy, 0.42, 0.42, INK)
        txt(sl, AX, yy + 0.09, 0.42, 0.24, '%02d' % (i + 1), 11, True, WHITE, PP_ALIGN.CENTER)
        txt(sl, AX + 0.58, yy + 0.01, 5.1, 0.24, a[:38], 11.5, True, INK)
        txt(sl, AX + 0.58, yy + 0.24, 5.1, 0.22, b[:52], 9.5, False, MUTED)
        yy += 0.62
    if alerts:
        rect(sl, AX, 5.62, 5.68, 1.18, RGBColor(0xFC, 0xF3, 0xF1), ACC)
        txt(sl, AX + 0.22, 5.72, 5.3, 0.22, '발표 전 확인 필요', 11.5, True, ACC)
        txt(sl, AX + 0.22, 5.96, 5.3, 0.78, '\n'.join('· ' + x[:58] for x in alerts[:3]),
            9.5, False, INK, line=1.25)
    foot(sl, src)


# ------------------------------------------------------------------ 조립
def card(ctx, item):
    l, p = ctx['last'].get(item, {}), ctx['prev'].get(item, {})
    t, tp = l.get('total', {}), p.get('total', {})
    pmap = {r['sty']: r for r in p.get('sty', [])}
    sty = []
    for r in l.get('sty', []):
        r = dict(r)
        r['_wow'] = wow(pmap.get(r['sty'], {}).get('t_gmv'), r.get('t_gmv'))
        sty.append(r)
    onw = wow(tp.get('on_gmv'), t.get('on_gmv'))
    offw = wow(tp.get('off_gmv'), t.get('off_gmv'))
    bullets = []
    if onw is not None or offw is not None:
        bullets.append('온라인 %s / 오프라인 %s' % (pct(onw), pct(offw)))
    bullets.append('목표비 %s (목표 %s개 / 실판 %s개)'
                   % (pct(t.get('t_ach'), False), ic(t.get('t_tgt')), ic(t.get('t_qty')))
                   if t.get('t_ach') else '주간 목표 미세팅 — 목표비 산출 불가')
    bullets.append('전년비 %s' % pct(wow(t.get('t_ly_gmv'), t.get('t_gmv')))
                   if t.get('t_ly_gmv') else '전년 동주 실적 없음(신상) — 전년비 해석 주의')
    msg = '온라인이 만든 성장' if (onw or 0) > (offw or 0) else '오프라인이 만든 성장'
    if (t.get('off_gmv') or 0) == 0:
        msg = '전량 온라인 채널'
    return dict(name=item, t=t, p=tp, sty=sty, psty=p.get('sty', []),
                wow=wow(tp.get('t_gmv'), t.get('t_gmv')),
                yoy=wow(t.get('t_ly_gmv'), t.get('t_gmv')) if t.get('t_ly_gmv') else None,
                ach=t.get('t_ach'), msg=msg, bullets=bullets)


def make_actions(ctx, cards):
    acts = []
    for c in cards:
        so, qo = g(c['t'], 'stock_on'), g(c['t'], 'on_qty')
        sf, qf = g(c['t'], 'stock_off'), g(c['t'], 'off_qty')
        if so and qo and sf and qf:
            w_on, w_off = so / qo, sf / qf
            if w_on < 12 and w_off > w_on * 3:
                acts.append(('%s 온라인 재고 보충' % c['name'],
                             '온 %.1f주 vs 오프 %.1f주 — 성장 채널이 먼저 마름' % (w_on, w_off)))
            elif w_off < 12 and w_on > w_off * 3:
                acts.append(('%s 오프라인 재고 보충' % c['name'],
                             '오프 %.1f주 vs 온 %.1f주' % (w_off, w_on)))
    for c in cards:
        big = [r for r in c['sty'] if (r.get('t_gmv') or 0) > 1e7 and (r['_wow'] or 0) > .5]
        for r in sorted(big, key=lambda r: -(r['_wow'] or 0))[:1]:
            onq, offq = r.get('on_qty') or 0, r.get('off_qty') or 0
            acts.append(('%s 온라인 확대' % (r.get('name') or '')[:20],
                         '전주비 %s인데 온라인 비중 %.0f%%'
                         % (pct(r['_wow']), onq / max(onq + offq, 1) * 100)))
    for item, inf in (ctx.get('inflow') or {}).items():
        dd = inf.get('daily') or []
        if len(dd) > 1:
            j = max(range(1, len(dd)), key=lambda i: dd[i][1] - dd[i - 1][1])
            if dd[j][1] > dd[j - 1][1] * 1.4:
                acts.append(('%s %s 유입 급점프 원인 규명' % (item, dd[j][0][5:]),
                             '%s → %s UV. 재현 가능한지 판단' % (ic(dd[j - 1][1]), ic(dd[j][1]))))
                break
    for c in cards:
        pend = [r for r in c['sty'] if r.get('pending')]
        if pend:
            acts.append(('%s 미발매 STY 투입 일정' % c['name'],
                         '%d종 대기 — %s 등' % (len(pend), (pend[0].get('name') or '')[:18])))
            break
    return acts


def collect_alerts(ctx):
    out = []
    for item in ctx['tabs']:
        t = ctx['last'].get(item, {}).get('total', {})
        st, so, sf = t.get('stock_t') or 0, t.get('stock_on') or 0, t.get('stock_off') or 0
        if st and abs(st - (so + sf)) > 1:
            out.append('%s 재고 합계 ≠ ON+OFF (%s개 미귀속)' % (item, ic(st - so - sf)))
        if so and sf and abs(so - sf) / max(so, sf) < 0.01:
            out.append('%s 온·오프 재고 동일값(%s) — 중복 계상 의심' % (item, ic(so)))
    if any(not ctx['last'].get(i, {}).get('total', {}).get('t_tgt') for i in ctx['deck_items']):
        out.append('주간 목표 미세팅 품목 있음 — 목표비 인용 주의')
    return out


def build(ctx, out_path):
    src = ('출처: 26FW 히어로 실적 대시보드(주간 블록) · 히어로 마스터앱 PDP 경로(direct)'
           ' | 실적 %s' % ctx['period'])
    cards = [card(ctx, i) for i in ctx['deck_items']]
    ctx['actions'] = make_actions(ctx, cards)
    alerts = ctx.get('alerts') or []
    d = Deck()
    s_cover(d, ctx, cards)
    s_summary(d, ctx, cards, src)
    s_split(d, ctx, cards[0], src)
    s_sty(d, ctx, cards[0], src)
    s_inflow(d, ctx, ctx['deck_items'][0], cards[0], src)
    if len(cards) > 1:
        s_waterfall(d, ctx, cards[1], src)
    s_action(d, ctx, cards, alerts, src)
    return d.save(out_path)


def upload(path, name):
    creds = get_credentials(WR._ROOT / 'credentials.json', WR._ROOT / 'token.json')
    drive = build_services(creds)['drive']
    f = drive.files().create(
        body={'name': name, 'mimeType': 'application/vnd.google-apps.presentation'},
        media_body=MediaFileUpload(
            path,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            resumable=True),
        fields='id,webViewLink').execute()
    return f['webViewLink']


def main():
    argv = sys.argv[1:]
    ctx = WR.main()
    if not ctx:
        raise SystemExit('시트 생성 결과가 없습니다.')
    ctx['alerts'] = collect_alerts(ctx)
    outdir = Path(os.environ.get('DECK_OUT') or tempfile.gettempdir())
    outdir.mkdir(parents=True, exist_ok=True)
    p = outdir / ('히어로_주간_세일즈_리뷰_%s.pptx' % ctx['week_end'].replace('-', ''))
    build(ctx, str(p))
    print('PPTX', p)
    if '--no-upload' not in argv:
        print('SLIDES', upload(str(p), '히어로 주간 세일즈 리뷰 (%s주)' % ctx['week_end']))
    print('SHEET', ctx['url'])
    return ctx


if __name__ == '__main__':
    main()
